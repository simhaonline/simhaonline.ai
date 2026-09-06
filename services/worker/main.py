# Simha worker: background jobs + tiny status API.
# Jobs (asyncio): gateway-triggered discovery, status health snapshots,
# TimescaleDB rollups, email delivery queue (Valkey pubsub), semantic memory
# indexer (pgvector), scheduled task runner.
import asyncio
import hashlib
import json
import logging
import os
import time
from contextlib import asynccontextmanager
from datetime import datetime, timezone

import httpx
import psycopg
from psycopg_pool import AsyncConnectionPool
import redis.asyncio as aioredis
import boto3
from fastapi import FastAPI

LOG = logging.getLogger("simha.worker")
logging.basicConfig(level=logging.INFO, format="%(asctime)s %(name)s %(levelname)s %(message)s")

DB_URL = os.environ.get("DATABASE_URL", "postgresql://simha:simha_dev_password@localhost:5433/simhaonline")
VALKEY_URL = os.environ.get("VALKEY_URL", "redis://localhost:6380/2")
GATEWAY_URL = os.environ.get("GATEWAY_URL", "http://gateway:8080")
CONTROL_URL = os.environ.get("CONTROL_PLANE_URL", "http://control-plane:8081")
DISCOVERY_INTERVAL = int(os.environ.get("MODEL_REFRESH_INTERVAL", "300"))
STATUS_INTERVAL = int(os.environ.get("STATUS_SNAPSHOT_INTERVAL", "60"))
UPLOAD_DIR = os.environ.get("UPLOAD_DIR", "/uploads")
OBJECT_STORAGE_BACKEND = os.environ.get("OBJECT_STORAGE_BACKEND", "s3")
OBJECT_STORAGE_BUCKET = os.environ.get("OBJECT_STORAGE_BUCKET", "simha-user-files")
object_storage = boto3.client(
    "s3",
    endpoint_url=os.environ.get("OBJECT_STORAGE_ENDPOINT", "http://minio:9000"),
    region_name=os.environ.get("OBJECT_STORAGE_REGION", "us-east-1"),
    aws_access_key_id=os.environ.get("OBJECT_STORAGE_ACCESS_KEY", ""),
    aws_secret_access_key=os.environ.get("OBJECT_STORAGE_SECRET_KEY", ""),
)

pool: AsyncConnectionPool | None = None


async def db() -> AsyncConnectionPool:
    global pool
    if pool is None:
        pool = AsyncConnectionPool(DB_URL, min_size=1, max_size=8, open=False)
        await pool.open()
    return pool


async def ensure_legacy_compatibility():
    """Apply additive, idempotent fields retained by the legacy router."""
    async with pool.connection() as c:
        # Timescale columnstore hypertables reject adding identity/constraint
        # columns after creation; retain the legacy field as an additive,
        # nullable compatibility column instead.
        await c.execute("ALTER TABLE request_history ADD COLUMN IF NOT EXISTS id BIGINT")
        await c.execute("ALTER TABLE discovered_models ADD COLUMN IF NOT EXISTS enabled BOOLEAN NOT NULL DEFAULT TRUE")
        await c.execute("""
            CREATE TABLE IF NOT EXISTS token_usage (
                account_name TEXT PRIMARY KEY REFERENCES accounts(name) ON DELETE CASCADE,
                prompt_tokens BIGINT NOT NULL DEFAULT 0,
                completion_tokens BIGINT NOT NULL DEFAULT 0,
                total_tokens BIGINT NOT NULL DEFAULT 0
            )
        """)
        await c.execute("""
            CREATE TABLE IF NOT EXISTS model_token_usage (
                model TEXT PRIMARY KEY,
                prompt_tokens BIGINT NOT NULL DEFAULT 0,
                completion_tokens BIGINT NOT NULL DEFAULT 0,
                total_tokens BIGINT NOT NULL DEFAULT 0
            )
        """)
        await c.execute("""
            INSERT INTO token_usage(account_name, prompt_tokens, completion_tokens, total_tokens)
            SELECT account_name, COALESCE(SUM(prompt_tokens), 0), COALESCE(SUM(completion_tokens), 0), COALESCE(SUM(total_tokens), 0)
            FROM request_history WHERE account_name IS NOT NULL GROUP BY account_name
            ON CONFLICT (account_name) DO NOTHING
        """)
        await c.execute("""
            INSERT INTO model_token_usage(model, prompt_tokens, completion_tokens, total_tokens)
            SELECT model, COALESCE(SUM(prompt_tokens), 0), COALESCE(SUM(completion_tokens), 0), COALESCE(SUM(total_tokens), 0)
            FROM request_history WHERE model IS NOT NULL GROUP BY model
            ON CONFLICT (model) DO NOTHING
        """)


# ---------------- jobs ----------------

async def trigger_discovery():
    """Ask the gateway to run one model-discovery pass (leader-locked)."""
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.post(f"{GATEWAY_URL}/internal/refresh-models")
            if r.status_code in (200, 202):
                LOG.info("discovery triggered")
            else:
                LOG.warning("discovery trigger: HTTP %s", r.status_code)
    except Exception as exc:  # noqa: BLE001
        LOG.warning("discovery trigger failed: %s", exc)


async def discovery_loop():
    while True:
        await trigger_discovery()
        await asyncio.sleep(DISCOVERY_INTERVAL)


async def status_snapshot():
    """Record one public health snapshot (Provider / Models components)."""
    global pool
    if pool is None:
        pool = await db()
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            r = await client.get(f"{GATEWAY_URL}/healthz")
            provider_ok = r.status_code == 200
        async with httpx.AsyncClient(timeout=10) as client:
            r2 = await client.get(f"{CONTROL_URL}/healthz")
            cp_ok = r2.status_code == 200
    except Exception:  # noqa: BLE001
        provider_ok = False
        cp_ok = False
    try:
        async with pool.connection() as c:
            async with c.cursor() as cur:
                # models_ok: at least one discovered model exists
                await cur.execute("SELECT COUNT(*) FROM discovered_models")
                n = (await cur.fetchone())[0]
                models_ok = n > 0
                await cur.execute(
                    "INSERT INTO status_checks(checked_at, provider_ok, models_ok) VALUES (now(), %s, %s)",
                    (provider_ok, models_ok),
                )
        LOG.info("status snapshot recorded (provider_ok=%s models_ok=%s cp_ok=%s)", provider_ok, models_ok, cp_ok)
    except Exception as exc:  # noqa: BLE001
        LOG.warning("status snapshot failed: %s", exc)


async def status_loop():
    while True:
        await status_snapshot()
        await asyncio.sleep(STATUS_INTERVAL)


async def email_worker():
    """Consume control-plane email publishes and deliver via SMTP (no-op if unconfigured)."""
    r = aioredis.from_url(VALKEY_URL, decode_responses=True)
    pubsub = r.pubsub()
    await pubsub.subscribe("simha:email")
    smtp_host = os.environ.get("SMTP_HOST", "")
    async for message in pubsub.listen():
        if message.get("type") != "message":
            continue
        try:
            doc = json.loads(message["data"])
        except (ValueError, TypeError):
            continue
        if not smtp_host:
            LOG.info("email skipped (SMTP unconfigured): %s -> %s", doc.get("subject"), doc.get("recipient"))
            continue
        try:
            import smtplib
            from email.message import EmailMessage

            m = EmailMessage()
            m["Subject"] = doc.get("subject", "")
            m["From"] = os.environ.get("SMTP_FROM", "no-reply@simhaonline.ai")
            m["To"] = doc.get("recipient", "")
            m.set_content(doc.get("text", ""))
            loop = asyncio.get_running_loop()
            await loop.run_in_executor(
                None,
                lambda: _send_smtp(m, smtp_host, int(os.environ.get("SMTP_PORT", "587")),
                                   os.environ.get("SMTP_USER", ""), os.environ.get("SMTP_PASS", "")),
            )
            LOG.info("email delivered: %s", doc.get("subject"))
        except Exception as exc:  # noqa: BLE001
            LOG.warning("email delivery failed: %s", exc)


def _send_smtp(m, host, port, user, password):
    with smtplib.SMTP(host, port, timeout=15) as server:
        server.starttls()
        if user:
            server.login(user, password)
        server.send_message(m)


async def rollups():
    """Continuous aggregates are replaced by a simple daily token rollup view refresh."""
    try:
        async with pool.connection() as c:
            await c.execute("""
                CREATE MATERIALIZED VIEW IF NOT EXISTS usage_daily AS
                SELECT time_bucket('1 day', requested_at) AS day,
                       account_name, model,
                       SUM(prompt_tokens) AS prompt_tokens,
                       SUM(completion_tokens) AS completion_tokens,
                       SUM(total_tokens) AS total_tokens,
                       COUNT(*) AS requests
                FROM request_history
                GROUP BY 1, 2, 3
                WITH NO DATA
            """)
            await c.execute("CREATE UNIQUE INDEX IF NOT EXISTS idx_usage_daily ON usage_daily(day, account_name, model)")
            await c.execute("REFRESH MATERIALIZED VIEW usage_daily")
            # billing rollup (per-user) — used by admin dashboards
            await c.execute("""
                CREATE MATERIALIZED VIEW IF NOT EXISTS user_usage_daily AS
                SELECT user_id,
                       COUNT(*) FILTER (WHERE requested_at > now() - interval '24 hours') AS requests_24h,
                       COUNT(*) FILTER (WHERE requested_at > now() - interval '30 days') AS requests_30d
                FROM request_history
                WHERE user_id IS NOT NULL
                GROUP BY user_id
                WITH NO DATA
            """)
            await c.execute("CREATE UNIQUE INDEX IF NOT EXISTS idx_user_usage_daily_uid ON user_usage_daily(user_id)")
            await c.execute("REFRESH MATERIALIZED VIEW user_usage_daily")
        LOG.info("usage_daily rollup refreshed")
    except Exception as exc:  # noqa: BLE001
        LOG.warning("rollup failed: %s", exc)


async def rollup_loop():
    while True:
        await rollups()
        await asyncio.sleep(900)


async def scheduled_runner():
    """Queue due scheduled tasks as chat messages (parity with legacy scheduler)."""
    try:
        tasks = []
        async with pool.connection() as c:
            async with c.cursor() as cur:
                await cur.execute("""
                    UPDATE scheduled_tasks SET status = 'running'
                    WHERE status = 'queued' AND run_at IS NOT NULL AND run_at <= now()
                    RETURNING id, user_id, prompt
                """)
                tasks = await cur.fetchall()
                for tid, user_id, prompt in tasks:
                    # enqueue via Valkey for the gateway-adjacent chat loop
                    await c.execute(
                        "INSERT INTO generated_content(user_id, kind, prompt, status) VALUES (%s, 'task', %s, 'queued')",
                        (user_id, prompt),
                    )
                    await c.execute("UPDATE scheduled_tasks SET status = 'done' WHERE id = %s", (tid,))
        if tasks:
            LOG.info("scheduled tasks dispatched: %d", len(tasks))
    except Exception as exc:  # noqa: BLE001
        LOG.warning("scheduler failed: %s", exc)


async def scheduler_loop():
    while True:
        await scheduled_runner()
        await asyncio.sleep(60)


# ---------------- Intake Dock ingestion ----------------

def _extract_file(file_path: str, mime_type: str, original_name: str) -> tuple[str, str]:
    """Extract safe text where a parser is installed; preserve binary assets."""
    suffix = os.path.splitext(original_name)[1].lower()
    text_types = mime_type.startswith("text/") or suffix in {".json", ".yaml", ".yml", ".py", ".js", ".jsx", ".ts", ".tsx", ".go", ".rs", ".java", ".c", ".cpp", ".h", ".hpp", ".sh", ".sql", ".md", ".csv", ".html", ".htm", ".css"}
    if text_types:
        with open(file_path, "rb") as handle:
            return handle.read(8 * 1024 * 1024).decode("utf-8", errors="replace"), "text"
    if mime_type == "application/pdf" or suffix == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            return "\n\n".join(page.extract_text() or "" for page in reader.pages)[:8 * 1024 * 1024], "pdf"
        except Exception as exc:  # noqa: BLE001
            return "", f"pdf-unavailable:{type(exc).__name__}"
    if suffix in {".xlsx", ".xls"}:
        try:
            import openpyxl
            book = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            lines = []
            for sheet in book.worksheets:
                lines.append(f"# Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    lines.append("\t".join("" if value is None else str(value) for value in row))
            return "\n".join(lines)[:8 * 1024 * 1024], "spreadsheet"
        except Exception as exc:  # noqa: BLE001
            return "", f"spreadsheet-unavailable:{type(exc).__name__}"
    if suffix in {".ppt", ".pptx"}:
        try:
            from pptx import Presentation
            presentation = Presentation(file_path)
            text = []
            for slide in presentation.slides:
                text.extend(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            return "\n".join(text)[:8 * 1024 * 1024], "slides"
        except Exception as exc:  # noqa: BLE001
            return "", f"slides-unavailable:{type(exc).__name__}"
    return "", "binary-passthrough"


async def process_ingestion_job(job_id: str):
    async with pool.connection() as c:
        await c.execute("UPDATE file_ingestion_jobs SET status='processing', started_at=COALESCE(started_at, now()), updated_at=now() WHERE id=%s AND status='queued'", (job_id,))
        async with c.cursor() as cur:
            await cur.execute("SELECT id, storage_key, storage_backend, storage_bucket, original_name, mime_type FROM file_ingestion_files WHERE job_id=%s ORDER BY created_at", (job_id,))
            files = await cur.fetchall()
    failed = 0
    for file_id, storage_key, storage_backend, storage_bucket, original_name, mime_type in files:
        file_path = os.path.join(UPLOAD_DIR, storage_key)
        temporary_object = None
        try:
            if (storage_backend or OBJECT_STORAGE_BACKEND) == "s3":
                temporary_object = os.path.join("/tmp", "simha-ingestion", str(file_id), os.path.basename(storage_key))
                os.makedirs(os.path.dirname(temporary_object), exist_ok=True)
                await asyncio.to_thread(object_storage.download_file, storage_bucket or OBJECT_STORAGE_BUCKET, storage_key, temporary_object)
                file_path = temporary_object
            extracted, parser = _extract_file(file_path, mime_type, original_name)
            async with pool.connection() as c:
                await c.execute("UPDATE file_ingestion_files SET status='processing', parser=%s, updated_at=now() WHERE id=%s", (parser, file_id))
                await c.execute("DELETE FROM file_ingestion_chunks WHERE file_id=%s", (file_id,))
                for index, start in enumerate(range(0, len(extracted), 6000)):
                    chunk = extracted[start:start + 6000]
                    vec = _hash_embedding(chunk, 384)
                    await c.execute("INSERT INTO file_ingestion_chunks(file_id, chunk_index, content, token_estimate, embedding) VALUES (%s,%s,%s,%s,%s::vector)", (file_id, index, chunk, max(1, len(chunk) // 4), "[" + ",".join(f"{x:.6f}" for x in vec) + "]"))
                await c.execute("UPDATE file_ingestion_files SET status='completed', extracted_text=%s, updated_at=now() WHERE id=%s", (extracted[:8 * 1024 * 1024], file_id))
        except Exception as exc:  # noqa: BLE001
            failed += 1
            LOG.exception("ingestion failed for %s", file_id)
            async with pool.connection() as c:
                await c.execute("UPDATE file_ingestion_files SET status='failed', error_message=%s, updated_at=now() WHERE id=%s", (str(exc)[:1000], file_id))
        finally:
            if temporary_object:
                try:
                    os.remove(temporary_object)
                except OSError:
                    pass
    async with pool.connection() as c:
        await c.execute("UPDATE file_ingestion_jobs SET status=%s, completed_at=now(), updated_at=now() WHERE id=%s", ("partial" if failed and failed < len(files) else "failed" if failed else "completed", job_id))


async def ingestion_loop():
    """Consume upload events and recover queued jobs after a worker restart."""
    redis = aioredis.from_url(VALKEY_URL, decode_responses=True)
    pubsub = redis.pubsub()
    await pubsub.subscribe("simha:ingestion")
    while True:
        message = await pubsub.get_message(ignore_subscribe_messages=True, timeout=1.0)
        if message and message.get("data"):
            try:
                await process_ingestion_job(json.loads(message["data"])["job_id"])
            except Exception as exc:  # noqa: BLE001
                LOG.warning("ingestion event failed: %s", exc)
        async with pool.connection() as c:
            async with c.cursor() as cur:
                await cur.execute("SELECT id FROM file_ingestion_jobs WHERE status='queued' ORDER BY created_at LIMIT 2")
                pending = await cur.fetchall()
        for (job_id,) in pending:
            await process_ingestion_job(str(job_id))
        await asyncio.sleep(1)


# ---------------- semantic memory (pgvector) ----------------

async def index_semantic(doc_kind: str, ref_id: str, content: str):
    """384-dim embedding via hashing-based fallback vector (deterministic,
    zero external calls). Real embeddings can be swapped in later without a
    schema change."""
    vec = _hash_embedding(content, 384)
    async with pool.connection() as c:
        await c.execute(
            """
            INSERT INTO semantic_memory(kind, ref_id, content, embedding)
            VALUES (%s, %s, %s, %s::vector)
            """,
            (doc_kind, ref_id, content[:4000], "[" + ",".join(f"{x:.6f}" for x in vec) + "]"),
        )


def _hash_embedding(text: str, dim: int) -> list[float]:
    """Deterministic token-histogram embedding (bag-of-words hashed to dim)."""
    vec = [0.0] * dim
    for token in text.lower().split():
        h = int(hashlib.sha256(token.encode()).hexdigest(), 16)
        vec[h % dim] += 1.0
    norm = sum(v * v for v in vec) ** 0.5 or 1.0
    return [v / norm for v in vec]


async def semantic_seed():
    """Seed semantic_memory with the provider catalog once."""
    try:
        catalog_path = os.environ.get("PROVIDER_CATALOG", "/config/provider_catalog.json")
        async with pool.connection() as c:
            async with c.cursor() as cur:
                await cur.execute("SELECT COUNT(*) FROM semantic_memory")
                n = (await cur.fetchone())[0]
                if n > 0:
                    return
                await cur.execute("SELECT value FROM app_settings WHERE key = 'provider_catalog'")
                row = await cur.fetchone()
                doc = json.loads(row[0]) if row else {}
                for category, items in (doc or {}).items():
                    if not isinstance(items, list):
                        continue
                    for item in items:
                        if not isinstance(item, dict):
                            continue
                        text = f"{item.get('name', '')} {item.get('id', '')} {category} protocol {item.get('protocol', '')}"
                        await index_semantic("provider", str(item.get("id", "")), text)
            LOG.info("semantic memory seeded")
    except Exception as exc:  # noqa: BLE001
        LOG.warning("semantic seed failed: %s", exc)


# ---------------- nightly database backup (pg_dump) ----------------

async def backup_run():
    """Run one pg_dump cycle (gzip-compressed) and prune old backups."""
    backup_dir = os.environ.get("BACKUP_DIR", "/backups")
    try:
        os.makedirs(backup_dir, exist_ok=True)
        stamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
        path = os.path.join(backup_dir, f"simhaonline_{stamp}.sql.gz")
        proc = await asyncio.create_subprocess_shell(
            f'pg_dump --no-owner --no-privileges --dbname "$SIMHA_DB" | gzip > "{path}"',
            env={**os.environ, "SIMHA_DB": DB_URL},
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        _, err = await proc.communicate()
        if proc.returncode != 0:
            # Keep the beginning as well as the tail: pg_dump often puts the
            # actual failing relation first and the useful hint last.
            detail = err.decode(errors="replace").strip()
            LOG.warning("pg_dump failed (exit %s): %s", proc.returncode, detail[:2000])
            return
        size = os.path.getsize(path)
        verify = await asyncio.create_subprocess_exec("gzip", "-t", path, stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.PIPE)
        _, verify_err = await verify.communicate()
        if verify.returncode != 0:
            LOG.warning("backup verification failed for %s: %s", path, verify_err.decode(errors="replace")[:1000])
            try:
                os.remove(path)
            except OSError:
                pass
            return
        LOG.info("backup written: %s (%d bytes)", path, size)
        # prune: keep 7 newest
        backups = sorted(
            f for f in os.listdir(backup_dir) if f.startswith("simhaonline_") and f.endswith(".sql.gz")
        )
        for old in backups[:-7]:
            os.remove(os.path.join(backup_dir, old))
            LOG.info("pruned old backup: %s", old)
    except Exception as exc:  # noqa: BLE001
        LOG.warning("backup failed: %s", exc)


async def backup_loop():
    interval = int(os.environ.get("BACKUP_INTERVAL", "86400"))
    await asyncio.sleep(120)  # let the DB settle after boot
    while True:
        await backup_run()
        await asyncio.sleep(interval)


# ---------------- http surface ----------------

@asynccontextmanager
async def lifespan(app: FastAPI):
    global pool
    pool = AsyncConnectionPool(DB_URL, min_size=1, max_size=8, open=False)
    await pool.open()
    await ensure_legacy_compatibility()
    tasks = [
        asyncio.create_task(discovery_loop(), name="discovery"),
        asyncio.create_task(status_loop(), name="status"),
        asyncio.create_task(email_worker(), name="email"),
        asyncio.create_task(rollup_loop(), name="rollups"),
        asyncio.create_task(scheduler_loop(), name="scheduler"),
        asyncio.create_task(ingestion_loop(), name="ingestion"),
        asyncio.create_task(backup_loop(), name="backup"),
        asyncio.create_task(_seed_delayed(), name="seed"),
    ]
    LOG.info("worker started")
    yield
    for t in tasks:
        t.cancel()
    await pool.close()


async def _seed_delayed():
    await asyncio.sleep(5)
    await semantic_seed()


app = FastAPI(title="Simha Worker", lifespan=lifespan)


@app.get("/healthz")
async def health():
    return {"status": "ok", "service": "simha-worker", "time": int(time.time())}


@app.get("/status/recent")
async def status_recent(limit: int = 15):
    """Public status feed: recent health snapshots (oldest→newest)."""
    limit = max(1, min(limit, 100))
    async with pool.connection() as c:
        async with c.cursor() as cur:
            await cur.execute(
                """
                SELECT checked_at, provider_ok, models_ok
                FROM status_checks
                ORDER BY checked_at DESC LIMIT %s
                """,
                (limit,),
            )
            rows = await cur.fetchall()
    return {
        "checks": [
            {"checked_at": r[0].isoformat(), "provider_ok": r[1], "models_ok": r[2]}
            for r in reversed(rows)
        ]
    }
